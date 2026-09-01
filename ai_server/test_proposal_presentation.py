"""PowerPoint 출력의 장수·핵심 섹션·기존 함수 계약을 확인합니다."""

from __future__ import annotations

import unittest

from pptx import Presentation

from ai_server.app.proposal_presentation import create_strategy_proposal_presentation


def _sample_report() -> dict:
    """네트워크와 OpenAI 없이 PPT 레이아웃을 검증할 최소 보고서를 만듭니다."""
    months = [f'2026{month:02d}' for month in range(1, 8)]
    return {
        'summary': '방문을 지역 소비와 체류로 연결하는 3개월 시범사업을 제안합니다.',
        'region_name': '서울특별시 강남구',
        'period': '2025-08~2026-07',
        'observed_findings': [
            {'metric': '월간 순 방문자 수', 'value': '17,963,441명', 'interpretation': '최신 월 관측값'},
            {'metric': '전년동월 증감률', 'value': '-8.3%', 'interpretation': '같은 달 대비'},
            {'metric': '관광소비 총액', 'value': '4,819억 원', 'interpretation': '외지인 관광소비'},
        ],
        'monthly_trend': [
            {'month': month, 'visitors': 17_000_000 + index * 120_000,
             'spending_krw': 450_000_000_000 + index * 5_000_000_000}
            for index, month in enumerate(months)
        ],
        'strategies': [{
            'timeframe': '3개월 시범 운영',
            'title': '야간 체류와 상권 소비를 연결하는 시범사업',
            'problem_to_solve': '방문 규모에 비해 숙박과 야간 소비 전환이 낮습니다.',
            'comparison_analysis': '같은 기준 비교지역보다 숙박 방문 비율이 낮은 편입니다.',
            'solution': '공식 관광지와 야간 상권을 연결하고 검증 가능한 혜택을 시범 운영합니다.',
            'expected_effect': '이용과 결제 변화를 같은 기준으로 확인해 확대 여부를 정할 수 있습니다.',
            'budget': '안내 페이지·운영 지원·측정 항목의 수량을 정하고 비교견적으로 산정합니다.',
            'kpi': '코스 이용 건수·관광소비액·숙박 방문 비율을 같은 기준으로 확인합니다.',
            'implementation_steps': [
                {'step': number, 'schedule': f'{number}단계', 'task': f'실행 작업 {number}',
                 'deliverable': f'확인 산출물 {number}'}
                for number in range(1, 6)
            ],
            'visual_asset_source_ids': [],
        }],
        'evidence_sources': [{
            'source_id': 'dataset:11680:1', 'source_type': 'dataset',
            'title': '한국관광 데이터랩 강남구 월간 원자료',
            'source_url': 'https://datalab.visitkorea.or.kr/',
        }],
        'execution_scenario': {'visitor_target_pct': 5, 'spending_target_pct': 8},
        'ml_analysis': {
            'status': 'available', 'model_version': 'test-v1', 'source_period': '202608~202610',
            'horizon_policy': {'schedule_status': 'undecided', 'decision_windows': []},
            'forecasts': [
                {'month': f'2026{month:02d}', 'visitors': 18_000_000 + index * 100_000,
                 'spending_krw': 470_000_000_000 + index * 4_000_000_000}
                for index, month in enumerate((8, 9, 10))
            ],
            'evaluation': {
                'validation_period': '202601~202603', 'test_period': '202604~202607',
                'metrics': {
                    'visitors': {'selected_model': 'RandomForestRegressor'},
                    'spending_krw': {'selected_model': 'seasonal_naive'},
                },
            },
        },
    }


class ProposalPresentationTest(unittest.TestCase):
    """PPT 다운로드가 서버에서 재현 가능한지 확인합니다."""

    def test_generates_eight_slide_strategy_deck(self) -> None:
        output = create_strategy_proposal_presentation(_sample_report())
        presentation = Presentation(output)
        self.assertEqual(len(presentation.slides), 8)
        all_text = '\n'.join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, 'text')
        )
        for label in ('핵심 판단과 제안', '문제와 기회', 'ML 자연추세', '5단계 실행 로드맵', '주요 공식 출처'):
            self.assertIn(label, all_text)


if __name__ == '__main__':
    unittest.main()
