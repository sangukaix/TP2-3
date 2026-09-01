"""검증된 관광 전략 보고서 JSON을 8장 실행기획서 PowerPoint로 만듭니다.

LLM이 만든 수치가 아니라 보고서에 이미 검증되어 들어온 관측값·ML 전망·사용자
목표만 사용합니다. 따라서 이 모듈은 새 분석이나 정책 효과 예측을 수행하지 않습니다.
"""

from __future__ import annotations

import re
from datetime import date
from io import BytesIO
from typing import Any
from urllib.parse import urlparse

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from .presentation_theme import (
    BLUE, BORDER, CYAN, INK, MUTED, NAVY, ORANGE, PALE_BLUE, PALE_CYAN,
    PALE_ORANGE, PURPLE, SLATE, SURFACE, WHITE, add_base_slide, add_bullets,
    add_chart_caption, add_image_cover, add_label, add_line, add_metric_card,
    add_photo_placeholder, add_rect, add_rich_text, add_series_chart, add_text,
    download_images, select_image_sources,
)


def _compact(value: Any, *, limit: int = 180) -> str:
    """보고서의 긴 문장을 슬라이드에 맞게 공백 정리하고 안전하게 줄입니다."""
    text = re.sub(r'\s+', ' ', str(value or '')).strip()
    if len(text) <= limit:
        return text
    shortened = text[:limit].rsplit(' ', 1)[0].rstrip(',. ')
    return f'{shortened}…'


def _sentences(value: Any, *, limit: int = 4, length: int = 90) -> list[str]:
    """장문을 표·목록에 넣기 좋은 짧은 문장으로 분리합니다."""
    text = re.sub(r'\s+', ' ', str(value or '')).strip()
    if not text:
        return []
    parts = [part.strip(' ·-') for part in re.split(r'(?<=[.!?다요])\s+|[;•]\s*', text) if part.strip()]
    if len(parts) == 1 and len(parts[0]) > length:
        parts = [part.strip() for part in re.split(r',\s*', parts[0]) if part.strip()]
    return [_compact(part, limit=length) for part in parts[:limit]]


def _as_float(value: Any) -> float:
    """차트 입력값의 None·문자열을 0으로 안전하게 바꿉니다."""
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _month_label(value: Any) -> str:
    """202609 또는 2026.09를 발표용 26.09 형식으로 바꿉니다."""
    digits = re.sub(r'\D', '', str(value or ''))
    if len(digits) >= 6:
        return f'{digits[2:4]}.{digits[4:6]}'
    return str(value or '')


def _format_people(value: float) -> str:
    """방문자 수를 만 명 단위로 간결하게 표시합니다."""
    return f'{value / 10_000:,.0f}만 명'


def _format_money(value: float) -> str:
    """관광소비액을 억 원 단위로 표시합니다."""
    return f'{value / 100_000_000:,.0f}억 원'


def _first_strategy(report: dict[str, Any]) -> dict[str, Any]:
    """현재 단일 전략 구조를 유지하되 빈 보고서에도 안전하게 동작합니다."""
    strategies = list(report.get('strategies') or [])
    return strategies[0] if strategies else {}


def _finding(report: dict[str, Any], index: int, fallback: str) -> tuple[str, str, str]:
    """상단 카드용 관측지표를 인덱스별로 읽습니다."""
    findings = list(report.get('observed_findings') or [])
    if index >= len(findings):
        return fallback, '자료 없음', ''
    item = findings[index]
    return (
        _compact(item.get('metric') or fallback, limit=34),
        _compact(item.get('value') or '자료 없음', limit=28),
        _compact(item.get('interpretation') or '', limit=58),
    )


def _brief_line(report: dict[str, Any]) -> str:
    """사업 여건을 표지에 한 줄로 표시하며 세부 입력은 본문에서 반복하지 않습니다."""
    brief = report.get('planning_brief') or {}
    pieces = []
    if brief.get('budget_status'):
        pieces.append(f"예산 {brief.get('budget_status')}")
    if brief.get('schedule_status'):
        pieces.append(f"일정 {brief.get('schedule_status')}")
    if brief.get('preferred_direction'):
        pieces.append(_compact(brief.get('preferred_direction'), limit=38))
    return '  ·  '.join(pieces[:3])


def _execution_targets(report: dict[str, Any]) -> tuple[float, float, bool]:
    """사용자가 화면에서 입력한 목표율만 읽습니다. 없으면 목표선을 만들지 않습니다."""
    scenario = report.get('execution_scenario') or {}
    if not scenario:
        return 0.0, 0.0, False
    return (
        max(0.0, _as_float(scenario.get('visitor_target_pct'))),
        max(0.0, _as_float(scenario.get('spending_target_pct'))),
        True,
    )


def _scenario_horizon(report: dict[str, Any], forecast_count: int) -> int:
    """미정 일정은 3개월, 확정된 희망 일정은 최대 6개월까지 유연하게 사용합니다."""
    ml = report.get('ml_analysis') or {}
    policy = ml.get('horizon_policy') or {}
    windows = list(policy.get('decision_windows') or [])
    if windows and str(policy.get('schedule_status') or '') not in {'', 'undecided', '미정'}:
        requested = int(_as_float(windows[0].get('months')) or 3)
        return max(1, min(requested, 6, forecast_count))
    return max(1, min(3, forecast_count))


def _scenario_rows(report: dict[str, Any]) -> dict[str, Any] | None:
    """ML 자연추세와 사용자가 정한 목표 도달 경로를 같은 월 기준으로 계산합니다.

    목표 경로는 인과효과 예측이 아니라, 목표를 달성하려면 필요한 월별 수준입니다.
    """
    ml = report.get('ml_analysis') or {}
    forecasts = list(ml.get('forecasts') or []) if ml.get('status') == 'available' else []
    if not forecasts:
        return None
    horizon = _scenario_horizon(report, len(forecasts))
    forecasts = forecasts[:horizon]
    visitor_target_pct, spending_target_pct, has_target = _execution_targets(report)
    categories = [_month_label(item.get('month')) for item in forecasts]
    baseline_visitors = [_as_float(item.get('visitors')) for item in forecasts]
    baseline_spending = [_as_float(item.get('spending_krw')) for item in forecasts]
    divisor = max(len(forecasts), 1)
    target_visitors = [
        value * (1 + visitor_target_pct / 100 * (index + 1) / divisor)
        for index, value in enumerate(baseline_visitors)
    ]
    target_spending = [
        value * (1 + spending_target_pct / 100 * (index + 1) / divisor)
        for index, value in enumerate(baseline_spending)
    ]
    return {
        'categories': categories,
        'baseline_visitors': baseline_visitors,
        'baseline_spending': baseline_spending,
        'target_visitors': target_visitors,
        'target_spending': target_spending,
        'visitor_target_pct': visitor_target_pct,
        'spending_target_pct': spending_target_pct,
        'visitor_gap': sum(target_visitors) - sum(baseline_visitors),
        'spending_gap': sum(target_spending) - sum(baseline_spending),
        'has_target': has_target,
        'horizon': horizon,
        'source_period': str(ml.get('source_period') or ''),
        'model_version': str(ml.get('model_version') or ''),
    }


def _model_summary(report: dict[str, Any]) -> tuple[str, str, str]:
    """방문자·소비액 모델과 검증기간을 출처 슬라이드용으로 압축합니다."""
    ml = report.get('ml_analysis') or {}
    evaluation = ml.get('evaluation') or {}
    metrics = evaluation.get('metrics') or {}
    visitor_model = str((metrics.get('visitors') or {}).get('selected_model') or '저장 모델')
    spending_model = str((metrics.get('spending_krw') or {}).get('selected_model') or '저장 모델')
    validation = str(evaluation.get('validation_period') or '')
    test = str(evaluation.get('test_period') or '')
    period = ' / '.join(item for item in (validation, test) if item) or '시간순 검증'
    return visitor_model, spending_model, period


def _add_cover(prs: Presentation, report: dict[str, Any], photo: tuple[dict[str, Any], BytesIO] | None) -> None:
    """큰 타이포와 지역 사진을 사용하는 표지입니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    if photo:
        source, stream = photo
        add_image_cover(slide, stream, 7.92, 0.02, 5.39, 7.46)
        add_rect(slide, 7.92, 0.02, 5.39, 7.46, fill=NAVY, line=NAVY, line_width=0, transparency=62)
        add_label(slide, _compact(source.get('title') or '공식 관광 이미지', limit=35), 8.25, 6.75, 4.55)
    else:
        add_photo_placeholder(slide, 7.92, 0.02, 5.39, 7.46)
    add_rect(slide, 0.58, 0.56, 0.72, 0.055, fill=BLUE, line=BLUE, line_width=0)
    add_text(slide, 'TOURISM STRATEGY PROPOSAL', 0.58, 0.75, 4.8, 0.22,
             size=8.5, color=BLUE, bold=True)
    add_text(slide, 'TOURISM', 0.42, 1.18, 7.1, 0.85, size=49, color=PALE_BLUE,
             bold=True, font_name='Arial')
    add_text(slide, str(report.get('region_name') or '선택 지역'), 0.62, 2.07, 6.65, 0.42,
             size=18, color=BLUE, bold=True)
    add_text(slide, '관광 전략\n실행 기획안', 0.62, 2.55, 6.78, 1.38, size=32, color=INK, bold=True)
    add_line(slide, 0.62, 4.22, 6.5, color=ORANGE, width=1.4)
    add_text(slide, _compact(report.get('summary'), limit=150), 0.62, 4.48, 6.48, 1.0,
             size=13, color=SLATE)
    brief = _brief_line(report)
    if brief:
        add_text(slide, brief, 0.62, 5.7, 6.5, 0.28, size=8.5, color=MUTED)
    add_text(slide, f"분석기간  {report.get('period') or '-'}", 0.62, 6.35, 3.8, 0.24,
             size=8.5, color=MUTED)
    add_text(slide, date.today().strftime('%Y.%m.%d'), 5.25, 6.35, 1.85, 0.24,
             size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, 'TOUR INSIGHT', 0.62, 6.83, 2.1, 0.22, size=8.5, color=NAVY, bold=True)


def _add_executive_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """서론을 늘리지 않고 판단·제안·목표를 한 장에 보여 줍니다."""
    strategy = _first_strategy(report)
    slide = add_base_slide(prs, '핵심 판단과 제안', 'EXECUTIVE PROPOSAL', 2)
    accents = [CYAN, BLUE, PURPLE]
    for index, fallback in enumerate(('월간 방문자', '관광소비액', '평균 숙박일수')):
        label, value, note = _finding(report, index, fallback)
        add_metric_card(slide, label, value, 0.58 + index * 4.08, 1.56, 3.82,
                        accent=accents[index], note=note)
    add_rect(slide, 0.58, 3.02, 5.77, 3.72, fill=WHITE, line=BORDER)
    add_text(slide, 'WHAT WE FOUND', 0.9, 3.32, 2.5, 0.2, size=8, color=ORANGE, bold=True)
    add_text(slide, '문제와 기회', 0.9, 3.63, 4.9, 0.38, size=20, color=INK, bold=True)
    diagnosis_items = _sentences(strategy.get('problem_to_solve'), limit=2, length=104)
    diagnosis_items += _sentences(strategy.get('comparison_analysis'), limit=2, length=104)
    add_bullets(slide, diagnosis_items, 0.9, 4.18, 4.98, 2.05, size=11.5,
                bullet_color=ORANGE, max_items=4)
    add_rect(slide, 6.62, 3.02, 6.13, 3.72, fill=NAVY, line=NAVY)
    add_text(slide, 'OUR PROPOSAL', 6.98, 3.32, 2.6, 0.2, size=8, color=CYAN, bold=True)
    add_text(slide, _compact(strategy.get('title') or '추천 전략', limit=58),
             6.98, 3.62, 5.34, 0.72, size=19, color=WHITE, bold=True)
    add_text(slide, _compact(strategy.get('solution'), limit=190),
             6.98, 4.47, 5.34, 1.18, size=12.5, color=WHITE)
    add_line(slide, 6.98, 5.87, 5.34, color=BLUE, width=1)
    add_text(slide, _compact(strategy.get('expected_effect'), limit=130),
             6.98, 6.08, 5.34, 0.45, size=9.5, color=PALE_CYAN)


def _actual_trend(report: dict[str, Any]) -> tuple[list[str], list[float], list[float]]:
    """관측값만 최근 12개월로 잘라 실제 추세 차트에 사용합니다."""
    rows = [item for item in (report.get('monthly_trend') or []) if not item.get('is_forecast')][-12:]
    return (
        [_month_label(item.get('month')) for item in rows],
        [_as_float(item.get('visitors')) / 10_000 for item in rows],
        [_as_float(item.get('spending_krw')) / 100_000_000 for item in rows],
    )


def _add_diagnosis_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """실제 원자료 추세와 타 지역 비교 근거를 같은 장에 배치합니다."""
    strategy = _first_strategy(report)
    categories, visitors, spending = _actual_trend(report)
    slide = add_base_slide(prs, '데이터에서 확인한 문제와 발전 가능성', 'DATA DIAGNOSIS', 3)
    add_rect(slide, 0.58, 1.55, 8.03, 5.45, fill=WHITE, line=BORDER)
    if categories:
        add_chart_caption(slide, '최근 12개월 방문자 추세', '단위: 만 명', 0.9, 1.86, 3.48)
        add_series_chart(
            slide, categories, [('방문자 수', visitors, CYAN)], 0.9, 2.2, 3.48, 2.05,
            chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED, number_format='#,##0', legend=False,
        )
        add_chart_caption(slide, '최근 12개월 관광소비 추세', '단위: 억 원', 4.74, 1.86, 3.48)
        add_series_chart(
            slide, categories, [('관광소비액', spending, BLUE)], 4.74, 2.2, 3.48, 2.05,
            number_format='#,##0', legend=False,
        )
    else:
        add_text(slide, '표시할 월별 관측값이 없습니다.', 0.9, 2.3, 7.3, 0.4,
                 size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_line(slide, 0.9, 4.56, 7.32, color=BORDER)
    add_text(slide, '비교 근거', 0.9, 4.84, 1.4, 0.22, size=9, color=BLUE, bold=True)
    add_text(slide, _compact(strategy.get('comparison_analysis'), limit=230),
             0.9, 5.2, 7.25, 1.17, size=12, color=SLATE)
    add_rect(slide, 8.87, 1.55, 3.88, 5.45, fill=PALE_ORANGE, line=ORANGE)
    add_text(slide, 'SO WHAT?', 9.22, 1.9, 1.8, 0.2, size=8.5, color=ORANGE, bold=True)
    add_text(slide, '지금 우선할 일', 9.22, 2.26, 3.08, 0.38, size=19, color=INK, bold=True)
    add_text(slide, _compact(strategy.get('problem_to_solve'), limit=235),
             9.22, 2.91, 3.12, 1.62, size=12.5, color=SLATE)
    add_line(slide, 9.22, 4.8, 3.12, color=ORANGE, width=1.3)
    add_text(slide, '관측 → 비교 → 제안', 9.22, 5.12, 3.08, 0.27, size=10.5, color=ORANGE, bold=True)
    add_text(slide, '원자료로 확인한 변화와 같은 기준의 지역 비교를 먼저 보고, 실행안은 그 다음에 제안합니다.',
             9.22, 5.58, 3.08, 0.93, size=10, color=SLATE)


def _add_strategy_slide(
    prs: Presentation,
    report: dict[str, Any],
    photo: tuple[dict[str, Any], BytesIO] | None,
) -> None:
    """하나의 실행 전략을 핵심 구조·운영 방식·산출물로 정리합니다."""
    strategy = _first_strategy(report)
    slide = add_base_slide(prs, '추천 전략을 한눈에 봅니다', 'STRATEGY DESIGN', 4)
    add_rect(slide, 0.58, 1.55, 5.02, 5.45, fill=NAVY, line=NAVY)
    add_text(slide, str(strategy.get('timeframe') or '3~6개월'), 0.92, 1.9, 2.25, 0.25,
             size=9, color=CYAN, bold=True)
    add_text(slide, _compact(strategy.get('title') or '추천 전략', limit=64),
             0.92, 2.32, 4.28, 1.02, size=22, color=WHITE, bold=True)
    add_text(slide, _compact(strategy.get('solution'), limit=260),
             0.92, 3.56, 4.27, 1.78, size=12.5, color=WHITE)
    add_line(slide, 0.92, 5.66, 4.27, color=BLUE, width=1)
    add_text(slide, '기대 변화', 0.92, 5.92, 1.2, 0.22, size=8.5, color=CYAN, bold=True)
    add_text(slide, _compact(strategy.get('expected_effect'), limit=120),
             0.92, 6.24, 4.27, 0.48, size=9.5, color=PALE_CYAN)
    if photo:
        source, stream = photo
        add_image_cover(slide, stream, 5.86, 1.55, 6.89, 3.22)
        add_label(slide, _compact(source.get('title') or '공식 관광 이미지', limit=42),
                  6.13, 4.15, 5.92, color=BLUE)
    else:
        add_photo_placeholder(slide, 5.86, 1.55, 6.89, 3.22)
    steps = list(strategy.get('implementation_steps') or [])[:3]
    card_width = 2.14
    for index in range(3):
        step = steps[index] if index < len(steps) else {}
        x = 5.86 + index * 2.36
        add_rect(slide, x, 5.02, card_width, 1.98, fill=WHITE, line=BORDER)
        add_text(slide, f'{index + 1:02d}', x + 0.18, 5.23, 0.46, 0.22,
                 size=9, color=BLUE, bold=True)
        add_text(slide, _compact(step.get('task') or '실행 단계', limit=58),
                 x + 0.18, 5.6, card_width - 0.36, 0.75, size=10.5, color=INK, bold=True)
        add_text(slide, _compact(step.get('deliverable') or '산출물', limit=38),
                 x + 0.18, 6.5, card_width - 0.36, 0.22, size=8, color=MUTED)


def _add_scenario_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """ML 자연추세와 사용자 실행 목표를 방문자·소비액 차트로 분리합니다."""
    scenario = _scenario_rows(report)
    slide = add_base_slide(prs, 'ML 자연추세와 실행 목표 시나리오', 'IMPACT SCENARIO', 5)
    if not scenario:
        add_rect(slide, 0.58, 1.56, 12.17, 4.72, fill=WHITE, line=BORDER)
        add_text(slide, '이 지역은 검증된 저장 ML 전망이 없어 목표 비교를 표시하지 않습니다.',
                 1.2, 3.12, 10.94, 0.5, size=17, color=MUTED, bold=True,
                 align=PP_ALIGN.CENTER)
        add_text(slide, '관측 자료와 공식 근거를 사용한 실행 기획은 나머지 장에서 계속 확인할 수 있습니다.',
                 1.2, 3.82, 10.94, 0.35, size=10.5, color=SLATE,
                 align=PP_ALIGN.CENTER)
        return
    categories = scenario['categories']
    add_rect(slide, 0.58, 1.55, 5.91, 3.83, fill=WHITE, line=BORDER)
    add_chart_caption(slide, '방문자 수: 자연추세 vs 실행 목표', '단위: 만 명', 0.9, 1.83, 5.27)
    visitor_series = [('ML 자연추세', [v / 10_000 for v in scenario['baseline_visitors']], BLUE)]
    if scenario['has_target']:
        visitor_series.append(('실행 목표', [v / 10_000 for v in scenario['target_visitors']], ORANGE))
    add_series_chart(slide, categories, visitor_series, 0.9, 2.18, 5.27, 2.72,
                     number_format='#,##0', legend=True)
    add_rect(slide, 6.84, 1.55, 5.91, 3.83, fill=WHITE, line=BORDER)
    add_chart_caption(slide, '관광소비액: 자연추세 vs 실행 목표', '단위: 억 원', 7.16, 1.83, 5.27)
    spending_series = [('ML 자연추세', [v / 100_000_000 for v in scenario['baseline_spending']], CYAN)]
    if scenario['has_target']:
        spending_series.append(('실행 목표', [v / 100_000_000 for v in scenario['target_spending']], ORANGE))
    add_series_chart(slide, categories, spending_series, 7.16, 2.18, 5.27, 2.72,
                     number_format='#,##0', legend=True)
    if scenario['has_target']:
        add_metric_card(
            slide, f"{scenario['horizon']}개월 누적 방문자 목표 차이",
            f"+{_format_people(scenario['visitor_gap'])}", 0.58, 5.68, 3.78,
            accent=ORANGE, note=f"최종월 목표 +{scenario['visitor_target_pct']:.1f}%",
        )
        add_metric_card(
            slide, f"{scenario['horizon']}개월 누적 소비 목표 차이",
            f"+{_format_money(scenario['spending_gap'])}", 4.56, 5.68, 3.78,
            accent=ORANGE, note=f"최종월 목표 +{scenario['spending_target_pct']:.1f}%",
        )
        add_rect(slide, 8.54, 5.68, 4.21, 1.12, fill=PALE_ORANGE, line=ORANGE)
        add_text(slide, '해석', 8.8, 5.91, 0.7, 0.2, size=8.5, color=ORANGE, bold=True)
        add_text(slide, '주황선은 정책효과 예측이 아니라 목표 달성에 필요한 월별 수준입니다.',
                 8.8, 6.21, 3.55, 0.38, size=9.2, color=SLATE)
    else:
        add_rect(slide, 0.58, 5.68, 12.17, 1.12, fill=PALE_BLUE, line=BLUE)
        add_text(slide, '실행 목표 미입력', 0.88, 5.94, 2.0, 0.25, size=10, color=BLUE, bold=True)
        add_text(slide, '화면에서 방문자·소비 목표율을 입력하면 자연추세와 필요한 목표 수준이 함께 표시됩니다.',
                 2.75, 5.94, 9.45, 0.28, size=10, color=SLATE)


def _budget_condition(report: dict[str, Any]) -> str:
    """사용자가 입력한 희망 예산을 전략 예산 산정 원칙과 함께 보여 줍니다."""
    brief = report.get('planning_brief') or {}
    budget = (
        brief.get('budget_amount')
        or brief.get('budget_range')
        or brief.get('budget_status')
        or '미정'
    )
    return _compact(budget, limit=54)


def _add_budget_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """예산과 KPI를 실제 집행 검토에 필요한 표 형태로 묶습니다."""
    strategy = _first_strategy(report)
    slide = add_base_slide(prs, '예산과 성과 확인 기준', 'BUDGET & KPI', 6)
    add_rect(slide, 0.58, 1.55, 4.0, 5.45, fill=NAVY, line=NAVY)
    add_text(slide, 'AVAILABLE BUDGET', 0.92, 1.9, 2.8, 0.2, size=8, color=CYAN, bold=True)
    add_text(slide, '희망 예산', 0.92, 2.33, 2.65, 0.34, size=18, color=WHITE, bold=True)
    add_text(slide, _budget_condition(report), 0.92, 2.9, 3.25, 0.55, size=22, color=CYAN, bold=True)
    add_line(slide, 0.92, 3.72, 3.28, color=BLUE, width=1)
    add_text(slide, '예산 산정 원칙', 0.92, 4.05, 2.0, 0.25, size=9, color=CYAN, bold=True)
    add_text(slide, _compact(strategy.get('budget'), limit=255),
             0.92, 4.46, 3.28, 1.54, size=11.5, color=WHITE)
    add_text(slide, '공식 단가 또는 항목×수량×단가 산식으로만 확정',
             0.92, 6.34, 3.28, 0.3, size=8.5, color=PALE_CYAN)
    add_rect(slide, 4.86, 1.55, 7.89, 2.5, fill=WHITE, line=BORDER)
    add_text(slide, 'BUDGET BREAKDOWN', 5.2, 1.87, 3.1, 0.2, size=8, color=BLUE, bold=True)
    budget_items = _sentences(strategy.get('budget'), limit=4, length=92)
    if not budget_items:
        budget_items = ['운영 항목과 수량을 정합니다.', '공식 단가 또는 2개 이상 비교견적을 확보합니다.']
    add_bullets(slide, budget_items, 5.2, 2.27, 7.15, 1.36, size=11,
                bullet_color=BLUE, max_items=4)
    add_rect(slide, 4.86, 4.3, 7.89, 2.7, fill=PALE_CYAN, line=CYAN)
    add_text(slide, 'MEASUREMENT', 5.2, 4.62, 2.45, 0.2, size=8, color=CYAN, bold=True)
    add_text(slide, '성과 확인 지표', 5.2, 4.96, 3.2, 0.34, size=18, color=INK, bold=True)
    kpi_items = _sentences(strategy.get('kpi'), limit=4, length=92)
    add_bullets(slide, kpi_items, 5.2, 5.46, 7.15, 1.15, size=11,
                bullet_color=CYAN, max_items=4)


def _add_timeline_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """5개 실행단계를 일정·행동·산출물 순서로 보여 줍니다."""
    strategy = _first_strategy(report)
    steps = list(strategy.get('implementation_steps') or [])[:5]
    slide = add_base_slide(prs, '5단계 실행 로드맵', 'ACTION TIMELINE', 7)
    add_text(slide, _compact(strategy.get('title'), limit=92), 0.58, 1.5, 9.65, 0.34,
             size=13, color=SLATE)
    add_text(slide, str(strategy.get('timeframe') or ''), 10.35, 1.51, 2.38, 0.28,
             size=10, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 1.18, 2.65, 11.0, color=CYAN, width=3)
    count = max(len(steps), 5)
    card_width = 2.28
    gap = 0.16
    left = 0.58
    for index in range(5):
        step = steps[index] if index < len(steps) else {}
        x = left + index * (card_width + gap)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(2.37), Inches(0.58), Inches(0.58),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE if index < count else MUTED
        circle.line.color.rgb = WHITE
        add_text(slide, str(step.get('step') or index + 1), x + 0.15, 2.37, 0.58, 0.58,
                 size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, _compact(step.get('schedule') or f'{index + 1}단계', limit=20),
                 x + 0.85, 2.47, card_width - 0.92, 0.22, size=8, color=BLUE, bold=True)
        add_rect(slide, x, 3.12, card_width, 3.45, fill=WHITE, line=BORDER)
        add_text(slide, _compact(step.get('task') or '실행 내용', limit=96),
                 x + 0.2, 3.43, card_width - 0.4, 1.25, size=11.3, color=INK, bold=True)
        add_line(slide, x + 0.2, 4.91, card_width - 0.4, color=BORDER)
        add_text(slide, '산출물', x + 0.2, 5.17, card_width - 0.4, 0.2,
                 size=8, color=MUTED, bold=True)
        add_text(slide, _compact(step.get('deliverable') or '확인 가능한 결과물', limit=58),
                 x + 0.2, 5.53, card_width - 0.4, 0.7, size=9.5, color=SLATE)


def _source_type_label(value: Any) -> str:
    """내부 source_type을 발표자가 이해하기 쉬운 출처명으로 바꿉니다."""
    return {
        'dataset': '관광데이터랩',
        'open_api': '관광 Open API',
        'tourism_open_api': '관광 Open API',
        'rag': '공식 문서',
        'web': '공식 웹 자료',
        'case_study': '공식 사례',
        'ml': '저장 ML 모델',
    }.get(str(value or '').lower(), '공식 자료')


def _source_domain(value: Any) -> str:
    """긴 URL 대신 도메인만 표시해 출처표를 읽기 쉽게 만듭니다."""
    try:
        return urlparse(str(value or '')).netloc or '-'
    except ValueError:
        return '-'


def _add_sources_slide(prs: Presentation, report: dict[str, Any]) -> None:
    """공식 자료와 ML 검증 방식을 마지막 장에서 짧게 확인합니다."""
    slide = add_base_slide(prs, '근거와 분석 방법', 'SOURCES & METHOD', 8)
    sources = list(report.get('evidence_sources') or [])[:6]
    add_rect(slide, 0.58, 1.55, 7.72, 5.46, fill=WHITE, line=BORDER)
    add_text(slide, '주요 공식 출처', 0.9, 1.84, 2.7, 0.3, size=17, color=INK, bold=True)
    add_text(slide, '자료명과 제공처만 간결하게 표시합니다.', 4.52, 1.9, 3.4, 0.18,
             size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)
    if not sources:
        add_text(slide, '한국관광 데이터랩 지역 원자료', 0.9, 2.55, 6.94, 0.35,
                 size=12, color=SLATE)
    for index, source in enumerate(sources):
        y = 2.34 + index * 0.68
        add_text(slide, f'{index + 1:02d}', 0.9, y + 0.03, 0.42, 0.2,
                 size=8.5, color=BLUE, bold=True)
        add_text(slide, _source_type_label(source.get('source_type')), 1.42, y, 1.35, 0.24,
                 size=8.5, color=BLUE, bold=True)
        add_text(slide, _compact(source.get('title') or '공식 자료', limit=65),
                 2.75, y, 3.7, 0.28, size=10.2, color=INK)
        add_text(slide, _source_domain(source.get('source_url')), 6.42, y + 0.02, 1.48, 0.22,
                 size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)
        add_line(slide, 0.9, y + 0.46, 6.98, color=BORDER, width=0.55)
    add_rect(slide, 8.56, 1.55, 4.19, 5.46, fill=NAVY, line=NAVY)
    visitor_model, spending_model, test_period = _model_summary(report)
    ml = report.get('ml_analysis') or {}
    ml_available = ml.get('status') == 'available'
    add_text(slide, 'ML METHOD', 8.92, 1.88, 2.4, 0.2, size=8, color=CYAN, bold=True)
    add_text(slide, '자연추세 전망', 8.92, 2.25, 3.2, 0.38, size=19, color=WHITE, bold=True)
    if ml_available:
        method_items = [
            f'방문자 모델: {visitor_model}',
            f'소비액 모델: {spending_model}',
            f'검증: {test_period}',
            '시간순 분리 + seasonal-naive 기준선 비교',
        ]
    else:
        method_items = ['검증된 저장 모델 없음', '관측 자료와 공식 출처만 사용']
    add_bullets(slide, method_items, 8.92, 2.91, 3.35, 1.9,
                size=10.5, color=WHITE, bullet_color=CYAN, max_items=4)
    add_line(slide, 8.92, 5.05, 3.34, color=BLUE, width=1)
    add_text(slide, '역할 구분', 8.92, 5.35, 1.2, 0.22, size=8.5, color=CYAN, bold=True)
    add_text(slide, 'MySQL=관측 사실\nML=자연추세 숫자\nRAG·공식 웹=문서 근거\nLLM=설명·전략',
             8.92, 5.75, 3.25, 0.92, size=10.3, color=WHITE)


def _validate_presentation(prs: Presentation) -> None:
    """다운로드 직전에 장수와 필수 섹션 누락을 코드로 막습니다."""
    if len(prs.slides) != 8:
        raise ValueError(f'PowerPoint는 표지 포함 정확히 8장이어야 합니다: {len(prs.slides)}장')
    all_text = '\n'.join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, 'text')
    )
    required = (
        '관광 전략',
        '실행 기획안',
        '문제와 기회',
        '추천 전략',
        '실행 목표 시나리오',
        '5단계 실행 로드맵',
        '주요 공식 출처',
    )
    missing = [label for label in required if label not in all_text]
    if missing:
        raise ValueError(f"PowerPoint 필수 섹션 누락: {', '.join(missing)}")


def create_strategy_proposal_presentation(report: dict[str, Any]) -> BytesIO:
    """현재 보고서 값으로 16:9, 8장 실행기획서 PPTX를 반환합니다."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 공식 관광 Open API 사진만 사용하며, 다운로드 실패 시 추상 배경으로 대체합니다.
    images = download_images(select_image_sources(report, limit=2))
    cover_photo = images[0] if images else None
    strategy_photo = images[1] if len(images) > 1 else cover_photo

    _add_cover(prs, report, cover_photo)
    _add_executive_slide(prs, report)
    _add_diagnosis_slide(prs, report)
    _add_strategy_slide(prs, report, strategy_photo)
    _add_scenario_slide(prs, report)
    _add_budget_slide(prs, report)
    _add_timeline_slide(prs, report)
    _add_sources_slide(prs, report)
    _validate_presentation(prs)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
