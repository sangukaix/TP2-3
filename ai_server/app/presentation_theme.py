"""관광 전략 PPT의 공통 디자인과 시각화 도구입니다.

슬라이드마다 좌표·색·폰트 코드를 반복하지 않도록 한곳에서 관리합니다.
보고서 내용은 ``proposal_presentation.py``가 결정하고, 이 파일은 표현만 담당합니다.
"""

from __future__ import annotations

from io import BytesIO
from typing import Any, Iterable
from urllib.parse import urlparse

import httpx
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


# 참고 템플릿의 네이비·시안·오렌지 조합을 프로젝트 색상으로 정리했습니다.
INK = RGBColor(24, 28, 35)
NAVY = RGBColor(15, 49, 91)
BLUE = RGBColor(0, 91, 176)
CYAN = RGBColor(0, 188, 210)
ORANGE = RGBColor(247, 101, 0)
PURPLE = RGBColor(105, 82, 185)
SLATE = RGBColor(74, 86, 104)
MUTED = RGBColor(139, 151, 168)
BORDER = RGBColor(216, 224, 234)
SURFACE = RGBColor(243, 247, 251)
PALE_BLUE = RGBColor(232, 244, 252)
PALE_CYAN = RGBColor(230, 249, 250)
PALE_ORANGE = RGBColor(255, 242, 232)
WHITE = RGBColor(255, 255, 255)


def set_run(run: Any, *, size: float, color: RGBColor = INK, bold: bool = False,
            font_name: str = '맑은 고딕') -> None:
    """텍스트 런의 글꼴 규칙을 통일합니다."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_text(
    slide: Any,
    value: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0,
    font_name: str = '맑은 고딕',
) -> Any:
    """인치 좌표로 텍스트 상자를 만들고 자동 줄바꿈을 켭니다."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    run.text = str(value or '')
    set_run(run, size=size, color=color, bold=bold, font_name=font_name)
    return box


def add_rich_text(
    slide: Any,
    segments: Iterable[tuple[str, bool, RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 15,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    """한 문장 안에서 숫자나 핵심어만 색·굵기를 달리 표시합니다."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, bold, color in segments:
        run = paragraph.add_run()
        run.text = text
        set_run(run, size=size, color=color, bold=bold)
    return box


def add_bullets(
    slide: Any,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: RGBColor = SLATE,
    bullet_color: RGBColor = CYAN,
    max_items: int = 5,
) -> Any:
    """긴 문단 대신 작은 색상 점과 짧은 문장을 배치합니다."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    clean_items = [str(item).strip() for item in items if str(item).strip()][:max_items]
    for index, item in enumerate(clean_items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(8)
        paragraph.line_spacing = 1.08
        run = paragraph.add_run()
        run.text = '● '
        set_run(run, size=max(7, size - 3), color=bullet_color, bold=True)
        run = paragraph.add_run()
        run.text = item
        set_run(run, size=size, color=color)
    return box


def add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = BORDER,
    radius: bool = False,
    line_width: float = 0.8,
    transparency: int = 0,
) -> Any:
    """카드·띠·배경으로 사용하는 사각형을 만듭니다."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    if line_width <= 0:
        # 전체 배경처럼 슬라이드 끝에 닿는 도형은 선을 제거해 렌더러 경계 초과를 막습니다.
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float = 0,
    *,
    color: RGBColor = BORDER,
    width: float = 1,
) -> Any:
    """섹션을 나누는 얇은 선을 추가합니다."""
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_base_slide(prs: Presentation, title: str, eyebrow: str, number: int) -> Any:
    """참고 템플릿과 같은 여백·상단선·큰 제목을 가진 기본 슬라이드입니다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SURFACE
    add_rect(slide, 0, 0, 13.333, 0.065, fill=BLUE, line=BLUE, line_width=0)
    add_text(slide, eyebrow.upper(), 0.58, 0.3, 4.8, 0.22, size=8.5, color=BLUE, bold=True)
    add_text(slide, title, 0.58, 0.62, 11.65, 0.55, size=25, color=INK, bold=True)
    add_text(slide, f'{number:02d}', 12.15, 0.33, 0.58, 0.26, size=9, color=MUTED,
             bold=True, align=PP_ALIGN.RIGHT)
    add_line(slide, 0.58, 1.28, 12.15, color=BORDER, width=0.8)
    return slide


def add_metric_card(
    slide: Any,
    label: str,
    value: str,
    x: float,
    y: float,
    w: float,
    *,
    accent: RGBColor = CYAN,
    note: str = '',
) -> None:
    """핵심 수치를 빠르게 읽는 얇은 카드입니다."""
    add_rect(slide, x, y, w, 1.12, fill=WHITE, line=BORDER)
    add_rect(slide, x, y, 0.055, 1.12, fill=accent, line=accent, line_width=0)
    add_text(slide, label, x + 0.2, y + 0.17, w - 0.38, 0.2, size=8.5, color=MUTED, bold=True)
    add_text(slide, value, x + 0.2, y + 0.43, w - 0.38, 0.38, size=18, color=INK, bold=True)
    if note:
        add_text(slide, note, x + 0.2, y + 0.84, w - 0.38, 0.17, size=7.5, color=SLATE)


def add_label(slide: Any, value: str, x: float, y: float, w: float, *, color: RGBColor = BLUE) -> None:
    """사진 캡션·상태 표시에 쓰는 작은 컬러 라벨입니다."""
    add_rect(slide, x, y, w, 0.32, fill=color, line=color, radius=False, line_width=0)
    add_text(slide, value, x + 0.09, y + 0.07, w - 0.18, 0.16, size=7.5, color=WHITE, bold=True)


def _valid_image_url(value: str) -> bool:
    """외부 사진은 HTTP(S) 주소만 허용합니다."""
    parsed = urlparse(str(value or ''))
    return parsed.scheme in {'http', 'https'} and bool(parsed.netloc)


def select_image_sources(report: dict[str, Any], *, limit: int = 2) -> list[dict[str, Any]]:
    """Agent가 선택한 관광 Open API 사진을 우선하고 나머지는 공식 사진으로 보완합니다."""
    sources = list(report.get('evidence_sources') or [])
    strategies = list(report.get('strategies') or [])
    selected_ids = {
        str(source_id)
        for strategy in strategies
        for source_id in (strategy.get('visual_asset_source_ids') or [])
    }
    candidates = [
        source for source in sources
        if _valid_image_url(str(source.get('image_url') or ''))
        and str(source.get('source_type') or '') in {'open_api', 'tourism_open_api'}
    ]
    candidates.sort(key=lambda item: 0 if str(item.get('source_id')) in selected_ids else 1)
    return candidates[:limit]


def download_images(sources: Iterable[dict[str, Any]]) -> list[tuple[dict[str, Any], BytesIO]]:
    """사진 다운로드 실패가 기획서 생성 전체를 막지 않도록 개별적으로 건너뜁니다."""
    downloaded: list[tuple[dict[str, Any], BytesIO]] = []
    with httpx.Client(timeout=7, follow_redirects=True) as client:
        for source in sources:
            try:
                response = client.get(str(source['image_url']))
                response.raise_for_status()
                if not str(response.headers.get('content-type', '')).lower().startswith('image/'):
                    continue
                stream = BytesIO(response.content)
                Image.open(stream).verify()
                stream.seek(0)
                downloaded.append((source, stream))
            except (httpx.HTTPError, OSError, KeyError, ValueError):
                continue
    return downloaded


def add_image_cover(slide: Any, stream: BytesIO, x: float, y: float, w: float, h: float) -> Any:
    """사진 비율을 유지하면서 지정한 프레임을 가득 채워 자릅니다."""
    stream.seek(0)
    with Image.open(stream) as image:
        image_ratio = image.width / max(image.height, 1)
    frame_ratio = w / h
    picture = slide.shapes.add_picture(stream, Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > frame_ratio:
        crop = max(0.0, (1 - frame_ratio / image_ratio) / 2)
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        crop = max(0.0, (1 - image_ratio / frame_ratio) / 2)
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def add_photo_placeholder(slide: Any, x: float, y: float, w: float, h: float) -> None:
    """공식 사진이 없을 때도 템플릿 균형을 유지하는 추상 배경입니다."""
    add_rect(slide, x, y, w, h, fill=NAVY, line=NAVY, line_width=0)
    add_rect(slide, x + w * 0.55, y, w * 0.45, h, fill=BLUE, line=BLUE,
             line_width=0, transparency=12)
    for index, alpha in enumerate((58, 68, 76, 84)):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.35 + index * 0.58),
            Inches(y + h - 1.65 - (index % 2) * 0.28),
            Inches(1.12 + index * 0.16),
            Inches(1.12 + index * 0.16),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = CYAN if index % 2 == 0 else WHITE
        circle.fill.transparency = alpha
        circle.line.fill.background()
    add_text(slide, 'REGIONAL\nTOURISM', x + 0.42, y + 0.42, w - 0.84, 1.2,
             size=26, color=WHITE, bold=True)


def add_series_chart(
    slide: Any,
    categories: list[str],
    series_rows: list[tuple[str, list[float], RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.LINE_MARKERS,
    number_format: str = '#,##0',
    legend: bool = True,
) -> Any:
    """실제 추세와 목표 시나리오에 공통으로 쓰는 최소형 차트입니다."""
    data = ChartData()
    data.categories = categories
    for name, values, _ in series_rows:
        data.add_series(name, values)
    chart = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), data).chart
    # python-pptx가 일부 막대차트 축 ID를 signed int로 만들 수 있습니다.
    # PowerPoint는 열지만 다른 OOXML 렌더러는 UInt32만 허용하므로 같은 참조값을 양수로 정규화합니다.
    for element in chart._chartSpace.xpath('.//*[@val]'):
        if not (element.tag.endswith('axId') or element.tag.endswith('crossAx')):
            continue
        raw_value = element.get('val')
        try:
            numeric_value = int(raw_value)
        except (TypeError, ValueError):
            continue
        if numeric_value < 0:
            element.set('val', str(numeric_value & 0xFFFFFFFF))
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.format.line.color.rgb = BORDER
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.number_format = number_format
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER
    chart.value_axis.format.line.color.rgb = BORDER
    for series, (_, _, color) in zip(chart.series, series_rows):
        if chart_type == XL_CHART_TYPE.LINE_MARKERS:
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2.25)
            series.marker.size = 6
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = WHITE
            series.marker.format.line.color.rgb = color
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.color.rgb = color
    return chart


def add_chart_caption(slide: Any, title: str, unit: str, x: float, y: float, w: float) -> None:
    """차트 제목과 단위를 한 줄에 정렬합니다."""
    add_text(slide, title, x, y, w * 0.72, 0.25, size=10.5, color=INK, bold=True)
    add_text(slide, unit, x + w * 0.72, y + 0.02, w * 0.28, 0.2, size=7.5,
             color=MUTED, align=PP_ALIGN.RIGHT)
