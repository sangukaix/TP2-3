"""검증된 관광 전략 보고서 JSON을 bid3 스타일의 PowerPoint 기획서로 만듭니다."""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
AQUA = RGBColor(14, 165, 168)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(148, 163, 184)
BORDER = RGBColor(226, 232, 240)
SURFACE = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


def _set_run(run, *, size: int, color: RGBColor = NAVY, bold: bool = False) -> None:
    run.font.name = '맑은 고딕'
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def _text(slide, text: str, x: float, y: float, w: float, h: float, *, size: int = 16,
          color: RGBColor = NAVY, bold: bool = False, align=PP_ALIGN.LEFT) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text or '')
    _set_run(run, size=size, color=color, bold=bold)
    return box


def _rect(slide, x: float, y: float, w: float, h: float, *, fill: RGBColor = WHITE,
          line: RGBColor = BORDER, radius: bool = True) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(.8)
    return shape


def _base_slide(prs: Presentation, title: str, section: str) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = SURFACE
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.08))
    top_line.fill.solid(); top_line.fill.fore_color.rgb = BLUE; top_line.line.color.rgb = BLUE
    _text(slide, section, .62, .34, 3.5, .25, size=9, color=BLUE, bold=True)
    _text(slide, title, .62, .58, 12.05, .55, size=23, bold=True)
    _text(slide, 'TOUR INSIGHT', 10.95, .35, 1.72, .25, size=8, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    return slide


def _add_metric_card(slide, label: str, value: str, x: float, accent: RGBColor) -> None:
    _rect(slide, x, 1.45, 3.75, 1.13)
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.45), Inches(.06), Inches(1.13))
    accent_line.fill.solid(); accent_line.fill.fore_color.rgb = accent; accent_line.line.color.rgb = accent
    _text(slide, label, x + .23, 1.66, 3.25, .25, size=9, color=SLATE)
    _text(slide, value, x + .23, 1.98, 3.25, .38, size=18, bold=True)


def _add_line_chart(slide, trend: list[dict[str, Any]], *, value_key: str, title: str,
                    x: float, y: float, w: float, h: float, color: RGBColor, divisor: float) -> None:
    chart_data = ChartData()
    chart_data.categories = [item.get('month', '') for item in trend]
    chart_data.add_series(title, [round(float(item.get(value_key) or 0) / divisor, 2) for item in trend])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h), chart_data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = BORDER
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(8)
    series = chart.series[0]
    series.format.line.color.rgb = color
    series.format.line.width = Pt(2.2)
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = WHITE
    series.marker.format.line.color.rgb = color


def create_strategy_proposal_presentation(report: dict[str, Any]) -> BytesIO:
    """화면과 동일한 보고서 값으로 5장 이내의 PPTX를 반환합니다."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    region = str(report.get('region_name') or '선택 지역')
    period = str(report.get('period') or '')
    trend = list(report.get('monthly_trend') or [])
    findings = list(report.get('observed_findings') or [])
    strategy = (report.get('strategies') or [{}])[0]
    sources = list(report.get('evidence_sources') or [])

    # 1. 표지
    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.background.fill.solid(); cover.background.fill.fore_color.rgb = NAVY
    cover_line = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.66), Inches(.75), Inches(.09), Inches(5.85))
    cover_line.fill.solid(); cover_line.fill.fore_color.rgb = AQUA; cover_line.line.color.rgb = AQUA
    _text(cover, 'REGIONAL TOURISM STRATEGY', 1.05, 1.25, 5.8, .35, size=10, color=RGBColor(94, 234, 212), bold=True)
    _text(cover, f'{region}\n관광 전략 기획안', 1.05, 1.75, 8.6, 1.5, size=31, color=WHITE, bold=True)
    _text(cover, str(report.get('summary') or ''), 1.05, 3.48, 8.4, 1.1, size=15, color=RGBColor(203, 213, 225))
    _text(cover, f'분석 기준  {period}', 1.05, 5.62, 4.4, .35, size=10, color=MUTED)
    _text(cover, 'TOUR INSIGHT', 10.15, 6.3, 2.1, .35, size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    # 2. 실제 관측값과 추세
    slide = _base_slide(prs, '현재 관광 현황', '01 · DATA SNAPSHOT')
    for index, finding in enumerate(findings[:3]):
        _add_metric_card(slide, str(finding.get('metric') or ''), str(finding.get('value') or ''), .62 + index * 4.04, [AQUA, BLUE, RGBColor(124, 58, 237)][index])
    if trend:
        _add_line_chart(slide, trend, value_key='visitors', title='월별 방문자 수 (백만 명)', x=.62, y=2.9, w=5.92, h=3.65, color=AQUA, divisor=1_000_000)
        _add_line_chart(slide, trend, value_key='spending_krw', title='월별 관광소비액 (십억 원)', x=6.78, y=2.9, w=5.92, h=3.65, color=BLUE, divisor=1_000_000_000)

    # 3. 문제/제안과 해결 방향
    slide = _base_slide(prs, '문제와 개선 방향', '02 · DIAGNOSIS')
    _rect(slide, .62, 1.45, 5.92, 4.9)
    _text(slide, '문제 / 발전 가능성', .92, 1.78, 5.3, .34, size=11, color=RGBColor(124, 58, 237), bold=True)
    _text(slide, str(strategy.get('problem_to_solve') or ''), .92, 2.22, 5.28, 1.12, size=16, bold=True)
    _text(slide, '이렇게 판단한 이유', .92, 3.6, 5.2, .3, size=10, color=SLATE, bold=True)
    _text(slide, str(strategy.get('comparison_analysis') or ''), .92, 3.98, 5.28, 1.58, size=12, color=SLATE)
    _rect(slide, 6.78, 1.45, 5.92, 4.9, fill=RGBColor(239, 246, 255), line=RGBColor(191, 219, 254))
    _text(slide, '추천 솔루션', 7.08, 1.78, 5.3, .34, size=11, color=BLUE, bold=True)
    _text(slide, str(strategy.get('title') or ''), 7.08, 2.22, 5.28, .78, size=19, bold=True)
    _text(slide, str(strategy.get('solution') or ''), 7.08, 3.12, 5.28, 1.6, size=13, color=SLATE)
    _text(slide, '기대 변화', 7.08, 4.95, 5.2, .28, size=10, color=AQUA, bold=True)
    _text(slide, str(strategy.get('expected_effect') or ''), 7.08, 5.28, 5.28, .7, size=11, color=SLATE)

    # 4. 실행 로드맵
    slide = _base_slide(prs, '실행 로드맵', '03 · ACTION PLAN')
    steps = list(strategy.get('implementation_steps') or [])[:5]
    if steps:
        left, total_width = .72, 11.9
        gap = .13
        card_width = (total_width - gap * (len(steps) - 1)) / len(steps)
        timeline_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(2.15), Inches(11.15), Inches(.035))
        timeline_line.fill.solid(); timeline_line.fill.fore_color.rgb = RGBColor(147, 197, 253); timeline_line.line.color.rgb = RGBColor(147, 197, 253)
        for index, step in enumerate(steps):
            x = left + index * (card_width + gap)
            _rect(slide, x, 1.55, card_width, 4.82)
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + .18), Inches(1.86), Inches(.46), Inches(.46))
            circle.fill.solid(); circle.fill.fore_color.rgb = BLUE; circle.line.color.rgb = BLUE
            _text(slide, str(step.get('step') or index + 1), x + .18, 1.86, .46, .46, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            _text(slide, str(step.get('schedule') or ''), x + .77, 1.88, card_width - .95, .28, size=9, color=BLUE, bold=True)
            _text(slide, str(step.get('task') or ''), x + .18, 2.62, card_width - .36, 1.45, size=12, bold=True)
            _text(slide, '완성되는 것', x + .18, 4.33, card_width - .36, .25, size=9, color=MUTED, bold=True)
            _text(slide, str(step.get('deliverable') or ''), x + .18, 4.65, card_width - .36, 1.12, size=10, color=SLATE)

    # 5. 검토 기준과 공식 출처
    slide = _base_slide(prs, '예산·성과 확인 기준', '04 · REVIEW')
    _rect(slide, .62, 1.45, 5.92, 2.1)
    _text(slide, '예산 산정 방법', .92, 1.78, 5.28, .3, size=11, color=BLUE, bold=True)
    _text(slide, str(strategy.get('budget') or ''), .92, 2.2, 5.28, 1.02, size=12, color=SLATE)
    _rect(slide, 6.78, 1.45, 5.92, 2.1)
    _text(slide, '성과 확인 지표', 7.08, 1.78, 5.28, .3, size=11, color=AQUA, bold=True)
    _text(slide, str(strategy.get('kpi') or ''), 7.08, 2.2, 5.28, 1.02, size=12, color=SLATE)
    _text(slide, '사용한 공식 자료', .62, 3.92, 4.2, .32, size=12, bold=True)
    source_lines = []
    for index, source in enumerate(sources[:6], start=1):
        source_lines.append(f"{index}. {source.get('title') or '공식 자료'}\n   {source.get('source_url') or ''}")
    if not source_lines:
        source_lines = ['관광데이터랩 지역 원자료 및 보고서에 표시된 공식 자료']
    _text(slide, '\n'.join(source_lines), .62, 4.33, 12.05, 1.75, size=9, color=SLATE)
    _text(slide, '※ 타 지역 사례의 성과는 선택 지역의 보장 효과가 아니며, 시범사업 전후 같은 기준으로 확인합니다.', .62, 6.52, 12.05, .3, size=8, color=MUTED)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
